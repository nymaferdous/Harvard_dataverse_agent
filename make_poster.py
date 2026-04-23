"""
make_poster.py
--------------
Generates MASBIO_Poster_2026.pptx in the style of the MASBio consortium poster:
two-column layout, dark-green branding, timeline, impact stats, dark footer.
Adds new Data Management section + Agent Architecture diagram.

Run:
    pip install python-pptx pillow requests
    python make_poster.py
"""

from pathlib import Path
import io, urllib.request
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Poster: A0 landscape (33.1 × 23.4 inches) ────────────────────────────────
W_IN, H_IN = 33.1, 23.4

def rgb(h):
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

# MASBio colour palette (matched from poster)
DARK_GREEN  = rgb("1B5E20")   # very dark green — logo text, headers
MED_GREEN   = rgb("2E7D32")   # medium green — subheadings
BRIGHT_GREEN= rgb("43A047")   # bright green — "Our Progress" style
LIGHT_GREEN = rgb("E8F5E9")   # very light green — section backgrounds
LIME        = rgb("76C442")   # MASBio lime accent
FOOTER_DARK = rgb("1A1A1A")   # near-black footer
RED         = rgb("C62828")   # accent red (annual meeting)
WHITE       = rgb("FFFFFF")
OFF_WHITE   = rgb("F9F9F9")
DARK_GRAY   = rgb("333333")
MID_GRAY    = rgb("666666")
LIGHT_GRAY  = rgb("F2F2F2")
TEAL        = rgb("006064")
MINT        = rgb("00897B")

FONT = "Calibri"

def in_(x): return Inches(x)
def pt_(x): return Pt(x)

# ── Core drawing helpers ──────────────────────────────────────────────────────

def rect(slide, x, y, w, h, fill, line_color=None, lw=0.5):
    from pptx.util import Pt
    sh = slide.shapes.add_shape(1, in_(x), in_(y), in_(w), in_(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line_color:
        sh.line.color.rgb = line_color; sh.line.width = Pt(lw)
    else:
        sh.line.fill.background()
    return sh

def txt(slide, text, x, y, w, h, size=10, bold=False, color=None,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(in_(x), in_(y), in_(w), in_(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.name = FONT; r.font.size = pt_(size)
    r.font.bold = bold; r.font.italic = italic
    if color: r.font.color.rgb = color
    return tb

def multirun(slide, runs, x, y, w, h, base_size=9, wrap=True, align=PP_ALIGN.LEFT):
    """runs = list of (text, bold, color, new_para)"""
    tb = slide.shapes.add_textbox(in_(x), in_(y), in_(w), in_(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    first = True
    for (text, bold, color, new_para) in runs:
        if new_para and not first:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        r = p.add_run(); r.text = text
        r.font.name = FONT; r.font.size = pt_(base_size)
        r.font.bold = bold
        if color: r.font.color.rgb = color
    return tb

def bullet_box(slide, items, x, y, w, h, size=9, color=None, bullet_color=None):
    """items = list of strings"""
    tb = slide.shapes.add_textbox(in_(x), in_(y), in_(w), in_(h))
    tf = tb.text_frame; tf.word_wrap = True
    color = color or DARK_GRAY
    bullet_color = bullet_color or MED_GREEN
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        # bullet symbol
        rb = p.add_run(); rb.text = "● "
        rb.font.name = FONT; rb.font.size = pt_(size)
        rb.font.color.rgb = bullet_color; rb.font.bold = True
        r = p.add_run(); r.text = item
        r.font.name = FONT; r.font.size = pt_(size)
        r.font.color.rgb = color

def stat_box(slide, number, label, x, y, size=1.4):
    """Impact statistic — green circle with number."""
    cx = x + size/2; cy = y + size/2
    # Circle
    circle = slide.shapes.add_shape(9, in_(x), in_(y), in_(size), in_(size))  # oval
    circle.fill.solid(); circle.fill.fore_color.rgb = LIGHT_GREEN
    circle.line.color.rgb = MED_GREEN; circle.line.width = Pt(2)
    # Number
    txt(slide, number, x, y+0.15, size, 0.65,
        size=20, bold=True, color=DARK_GREEN, align=PP_ALIGN.CENTER)
    # Label
    txt(slide, label, x-0.1, y+size-0.5, size+0.2, 0.55,
        size=9, color=DARK_GRAY, align=PP_ALIGN.CENTER)

def flow_node(slide, x, y, w, h, label, sublabel, fill, tc=WHITE):
    rect(slide, x, y, w, h, fill)
    tb = slide.shapes.add_textbox(in_(x+0.05), in_(y), in_(w-0.1), in_(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.name = FONT; r.font.size = pt_(8); r.font.bold = True; r.font.color.rgb = tc
    if sublabel:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = sublabel
        r2.font.name = FONT; r2.font.size = pt_(6.5); r2.font.color.rgb = tc

def section_header(slide, text, x, y, w, color=MED_GREEN, size=13):
    txt(slide, text, x, y, w, 0.4, size=size, bold=True, color=color)
    # underline rule
    rect(slide, x, y+0.38, w, 0.04, color)

def img_placeholder(slide, x, y, w, h, label, fill=LIGHT_GREEN, tc=MED_GREEN):
    rect(slide, x, y, w, h, fill, MED_GREEN, 0.5)
    txt(slide, label, x, y + h/2 - 0.2, w, 0.4,
        size=8, color=tc, align=PP_ALIGN.CENTER, italic=True)

# ════════════════════════════════════════════════════════════════════════════
# BUILD
# ════════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = in_(W_IN)
prs.slide_height = in_(H_IN)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# ── BACKGROUND ────────────────────────────────────────────────────────────────
rect(slide, 0, 0, W_IN, H_IN, WHITE)

# ════════════════════════════════════════════════════════════════════════════
# HEADER
# ════════════════════════════════════════════════════════════════════════════
HEADER_H = 2.6
rect(slide, 0, 0, W_IN, HEADER_H, WHITE)
rect(slide, 0, HEADER_H, W_IN, 0.07, DARK_GREEN)  # green rule under header

# MASBio logo text (left side)
txt(slide, "MASBio", 0.35, 0.12, 6.5, 1.5,
    size=72, bold=True, color=DARK_GREEN, align=PP_ALIGN.LEFT)
txt(slide, "DATA MANAGEMENT & PUBLISHING  |  2026",
    0.35, 1.55, 9.0, 0.5,
    size=13, bold=True, color=MED_GREEN)
txt(slide, "THE MID-ATLANTIC SUSTAINABLE BIOMASS FOR VALUE-ADDED PRODUCTS CONSORTIUM",
    0.35, 2.05, 12.0, 0.45,
    size=9.5, bold=False, color=DARK_GRAY)

# Authors
txt(slide,
    "Syeda Nyma Ferdous¹  ·  Karen Lopez-Olmedo¹  ·  Jinxing Wang¹  ·  Daniel Ciolkosz²  ·  Jude Liu²",
    0.35, 2.3, 14.0, 0.35, size=8.5, color=MID_GRAY)
txt(slide, "1. North Carolina State University   2. The Pennsylvania State University",
    0.35, 2.55, 10.0, 0.3, size=8, color=MID_GRAY, italic=True)

# Partner logos area (right side of header)
rect(slide, 18.5, 0.2, 14.3, 2.2, LIGHT_GRAY, MED_GREEN, 0.3)
txt(slide, "Partner Institutions",
    18.6, 0.22, 6.0, 0.35, size=8.5, bold=True, color=MED_GREEN)
# University name pills
unis = [
    ("West Virginia University", 18.6, 0.62),
    ("NC State University",       22.5, 0.62),
    ("Penn State",                26.0, 0.62),
    ("Virginia Tech",             29.2, 0.62),
    ("University of Wisconsin",   18.6, 1.1),
    ("USDA Forest Service",       22.5, 1.1),
    ("University of Georgia",     26.0, 1.1),
    ("Oak Ridge National Lab",    29.2, 1.1),
    ("West Virginia State Univ.", 18.6, 1.58),
    ("Dominion Energy",           22.5, 1.58),
    ("CTI Energy",                26.0, 1.58),
    ("AllStar Ecology",           29.2, 1.58),
]
for (name, ux, uy) in unis:
    rect(slide, ux, uy, 3.3, 0.38, WHITE, MED_GREEN, 0.4)
    txt(slide, name, ux+0.05, uy+0.05, 3.2, 0.3,
        size=7.5, bold=False, color=DARK_GREEN, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# CONTENT AREA
# ════════════════════════════════════════════════════════════════════════════
TOP   = HEADER_H + 0.18
BOT   = H_IN - 1.55
GAP   = 0.3
C1X   = 0.3
C1W   = 10.2
C2X   = C1X + C1W + GAP
C2W   = 11.0
C3X   = C2X + C2W + GAP
C3W   = W_IN - C3X - 0.3

# ════════════════════════════════════════════════════════════════════════════
# COLUMN 1 — About + Data Management
# ════════════════════════════════════════════════════════════════════════════

# About MASBio
cy = TOP
section_header(slide, "About MASBio", C1X, cy, C1W)
cy += 0.5
multirun(slide, [
    ("MASBio", True, DARK_GREEN, False),
    (" is a consortium of scientists from academia, government agencies and partners "
     "from the wood product industry. We collaborate on integrated and transdisciplinary "
     "research, education and extension to facilitate development of the bioeconomy and "
     "rural prosperity in the Mid-Atlantic region across seven primary areas:", False, DARK_GRAY, False),
], C1X, cy, C1W, 1.1, base_size=9.5)
cy += 1.15

# Research areas — two columns of bullets
areas_L = ["Feedstock Production", "Harvest and Logistics", "Bioproduct Production", "Sustainability Analysis"]
areas_R = ["System TEA and Optimization", "Education", "Outreach and Business"]
for i, a in enumerate(areas_L):
    rect(slide, C1X, cy + i*0.42, 0.3, 0.3, MED_GREEN)
    txt(slide, a, C1X+0.35, cy + i*0.42, C1W/2-0.4, 0.35,
        size=9, bold=True, color=DARK_GREEN)
for i, a in enumerate(areas_R):
    rect(slide, C1X + C1W/2, cy + i*0.42, 0.3, 0.3, LIME)
    txt(slide, a, C1X + C1W/2 + 0.35, cy + i*0.42, C1W/2-0.4, 0.35,
        size=9, bold=True, color=DARK_GREEN)
cy += 4*0.42 + 0.25

# Photos placeholder strip
img_placeholder(slide, C1X,       cy, 3.3, 2.4, "Field Research\nPhoto", LIGHT_GREEN)
img_placeholder(slide, C1X+3.4,   cy, 3.3, 2.4, "Team Meeting\nPhoto",   rgb("FFF9C4"), rgb("F9A825"))
img_placeholder(slide, C1X+6.8,   cy, 3.3, 2.4, "Biomass Sample\nPhoto", rgb("E3F2FD"), rgb("1565C0"))
cy += 2.55

# ── DATA MANAGEMENT SECTION (highlighted) ────────────────────────────────────
rect(slide, C1X, cy, C1W, 5.35, LIGHT_GREEN, DARK_GREEN, 1.0)
rect(slide, C1X, cy, C1W, 0.52, DARK_GREEN)  # header bar
txt(slide, "📂  DATA MANAGEMENT ACHIEVEMENTS",
    C1X+0.15, cy+0.07, C1W-0.25, 0.42,
    size=12, bold=True, color=WHITE)
cy += 0.6

# Sub-section: Folder reorganization
txt(slide, "Organized MASBIO Data Folder Structure",
    C1X+0.2, cy, C1W-0.3, 0.38,
    size=10.5, bold=True, color=DARK_GREEN)
cy += 0.42
bullet_box(slide, [
    "Reorganized all MASBIO data files into a standardized folder hierarchy for easy access by all consortium members",
    "Folders structured by project year, data type (sensor, production, policy), and institution",
    "Consistent naming conventions applied across all files and directories",
    "Shared drive permissions updated — all partners have read/write access to relevant folders",
], C1X+0.2, cy, C1W-0.35, 1.85, size=9, color=DARK_GRAY, bullet_color=MED_GREEN)
cy += 1.95

# Sub-section: Harvard Dataverse
txt(slide, "Harvard Dataverse — Open Data Repository",
    C1X+0.2, cy, C1W-0.3, 0.38,
    size=10.5, bold=True, color=DARK_GREEN)
cy += 0.42
bullet_box(slide, [
    "All processed MASBIO datasets published to Harvard Dataverse for long-term open access",
    "Each dataset assigned a unique, permanent Digital Object Identifier (DOI) — citable in publications",
    "Metadata auto-generated using AI pipeline: Groq LLM + RAG knowledge retrieval",
    "Datasets licensed under CC0 1.0 — fully open, no restrictions on reuse",
    "Compliant with FAIR data principles: Findable, Accessible, Interoperable, Reusable",
], C1X+0.2, cy, C1W-0.35, 2.1, size=9, color=DARK_GRAY, bullet_color=MED_GREEN)
cy += 2.15

# Harvard Dataverse badge
rect(slide, C1X+0.2, cy-0.15, C1W-0.4, 0.45, WHITE, MED_GREEN, 0.5)
txt(slide, "🔗  dataverse.harvard.edu  ·  DOI: 10.7910/DVN/MASBIO",
    C1X+0.3, cy-0.12, C1W-0.5, 0.38,
    size=9, bold=True, color=DARK_GREEN, align=PP_ALIGN.CENTER)

cy = TOP  # reset for column 2

# ════════════════════════════════════════════════════════════════════════════
# COLUMN 2 — Agent Architecture + Progress Timeline
# ════════════════════════════════════════════════════════════════════════════

section_header(slide, "AI-Powered Data Publishing Agent", C2X, cy, C2W, DARK_GREEN, 13)
cy += 0.5

txt(slide,
    "An autonomous LangChain agent automates the full Harvard Dataverse publishing workflow — "
    "from CSV upload to permanent DOI — using a hybrid Groq LLM + RAG metadata pipeline.",
    C2X, cy, C2W, 0.65, size=9.5, color=DARK_GRAY)
cy += 0.72

# ── Agent diagram ────────────────────────────────────────────────────────────
DX = C2X + 0.15
DW = C2W - 0.3
bw = 3.8; bh = 0.6
cx_ = DX + (DW-bw)/2

def ar(slide, cx, y1, y2):
    sh = slide.shapes.add_connector(1, in_(cx), in_(y1), in_(cx), in_(y2))
    sh.line.color.rgb = MED_GREEN; sh.line.width = Pt(1.5)

# Row 1
flow_node(slide, cx_, cy, bw, bh, "Researcher uploads CSV", "Browser UI — no code required", DARK_GREEN)
ar(slide, cx_+bw/2, cy+bh, cy+bh+0.22)

cy += bh + 0.22
# Row 2: RAG + Groq side by side
hw = (DW-0.2)/2
flow_node(slide, DX, cy, hw, 0.85,
    "RAG — Local Feature Extraction",
    "Time period · Geography · Data source\nsentence-transformers (free, offline)", TEAL)
flow_node(slide, DX+hw+0.15, cy, hw, 0.85,
    "Groq LLM — Llama 3 70B",
    "Description · Keywords · Subject\nFree API — no credits needed", MINT)

# merge arrows
ar(slide, DX+hw/2,         cy+0.85, cy+0.85+0.2)
ar(slide, DX+hw+0.15+hw/2, cy+0.85, cy+0.85+0.2)
cy += 0.85 + 0.2

# Merge row
mw = 5.5
flow_node(slide, DX+(DW-mw)/2, cy, mw, 0.5,
    "Merged Metadata  —  Title always from filename",
    None, MED_GREEN)
ar(slide, cx_+bw/2, cy+0.5, cy+0.5+0.2)
cy += 0.5 + 0.2

# Agent
aw = DW - 0.1
flow_node(slide, DX+0.05, cy, aw, 0.55,
    "LangChain ReAct Agent",
    "Reasons step-by-step · Selects tools · Handles errors autonomously",
    DARK_GREEN)
cy += 0.55 + 0.18

# 5 tools
tools = [
    ("Search", "Duplicates"),
    ("Create", "POST → DOI"),
    ("Upload", "CSV → Draft"),
    ("Verify", "Check state"),
    ("Publish", "Lock → Live"),
]
tcolors = [TEAL, MED_GREEN, MINT, rgb("006064"), DARK_GREEN]
tw_ = (aw - 0.04*4) / 5
for i,(l,s) in enumerate(tools):
    tx_ = DX + 0.05 + i*(tw_+0.04)
    ar(slide, tx_+tw_/2, cy-0.18, cy)
    flow_node(slide, tx_, cy, tw_, 0.72, l, s, tcolors[i])
ar(slide, cx_+bw/2, cy+0.72, cy+0.72+0.2)
cy += 0.72 + 0.2

# Harvard Dataverse
dvw = 4.5
flow_node(slide, DX+(DW-dvw)/2, cy, dvw, 0.6,
    "Harvard Dataverse",
    "Permanent DOI · CC0 1.0 · Publicly accessible",
    rgb("8B0000"))
cy += 0.6 + 0.25

# Result
txt(slide, "✅  Dataset published — DOI returned to browser in < 60 seconds",
    C2X, cy, C2W, 0.4, size=9.5, bold=True, color=DARK_GREEN, align=PP_ALIGN.CENTER)
cy += 0.55

# ── Progress Timeline ─────────────────────────────────────────────────────────
section_header(slide, "Our Progress", C2X, cy, C2W, BRIGHT_GREEN, 14)
cy += 0.5

# Timeline bar
rect(slide, C2X+0.3, cy+0.35, C2W-0.6, 0.08, BRIGHT_GREEN)
years   = ["2020–21", "2022", "2023", "2024", "2025", "2026"]
mile_x  = [C2X+0.3 + i*(C2W-0.6)/5 for i in range(6)]
events  = [
    "Project\nstarted",
    "Field trials\ndata collected",
    "Biochar demos\nNew bioecon. courses",
    "20+ publications\nML biochar model",
    "5th Annual\nMeeting NCSU",
    "Data Mgmt\n& Closure",
]
for i,(mx,yr,ev) in enumerate(zip(mile_x, years, events)):
    # dot
    dot = slide.shapes.add_shape(9, in_(mx-0.12), in_(cy+0.22), in_(0.24), in_(0.24))
    dot.fill.solid(); dot.fill.fore_color.rgb = DARK_GREEN
    dot.line.fill.background()
    # year
    txt(slide, yr, mx-0.45, cy, 0.9, 0.28,
        size=8.5, bold=True, color=DARK_GREEN, align=PP_ALIGN.CENTER)
    # event
    txt(slide, ev, mx-0.6, cy+0.5, 1.2, 0.6,
        size=7.5, color=DARK_GRAY, align=PP_ALIGN.CENTER)
cy += 1.3

# DATA management milestone highlight
rect(slide, C2X, cy, C2W, 0.55, LIGHT_GREEN, DARK_GREEN, 0.5)
txt(slide,
    "2025–2026:  MASBIO Data Folder reorganized  ·  All datasets published to Harvard Dataverse  ·  Unique DOI per file",
    C2X+0.15, cy+0.1, C2W-0.25, 0.38,
    size=9, bold=True, color=DARK_GREEN, align=PP_ALIGN.CENTER)
cy += 0.65

# Image placeholders (field work style)
img_placeholder(slide, C2X,        cy, 3.5, 2.0, "Biomass Field\nTrials Photo",       LIGHT_GREEN)
img_placeholder(slide, C2X+3.6,    cy, 3.5, 2.0, "Biochar Research\nLab Photo",       rgb("E3F2FD"), rgb("1565C0"))
img_placeholder(slide, C2X+7.2,    cy, 3.5, 2.0, "Harvard Dataverse\nDataset Page",   rgb("FFF8E1"), rgb("F57F17"))
cy += 2.1

# ════════════════════════════════════════════════════════════════════════════
# COLUMN 3 — Impact Stats + Annual Meeting
# ════════════════════════════════════════════════════════════════════════════

cy3 = TOP
section_header(slide, "Project Impact", C3X, cy3, C3W, DARK_GREEN, 13)
cy3 += 0.55

# 6 stat circles in 2 rows
stats = [
    ("+1000", "People\nEngaged"),
    ("+59",   "Students\n(UG, MS, PhD)"),
    ("+40",   "Outreach\nActivities"),
    ("+15",   "Programs\nDelivered"),
    ("+200",  "Academic\nPresentations"),
    ("+49",   "Peer-Reviewed\nPublications"),
]
scols = 3; sw = 1.55
for i,(num,lab) in enumerate(stats):
    row = i // scols; col = i % scols
    sx = C3X + col*(sw+0.28)
    sy = cy3 + row*(sw+0.55)
    stat_box(slide, num, lab, sx, sy, sw)
cy3 += 2*(sw+0.55) + 0.2

# ── Data Management Stats ─────────────────────────────────────────────────────
section_header(slide, "Data Management Impact", C3X, cy3, C3W, DARK_GREEN, 13)
cy3 += 0.55
dm_stats = [
    ("100%", "Datasets with\nUnique DOI"),
    ("CC0",  "Open License\non all datasets"),
    ("< 60s","Publish time\nper dataset"),
]
for i,(num,lab) in enumerate(dm_stats):
    sx = C3X + i*(sw+0.28)
    stat_box(slide, num, lab, sx, cy3, sw)
cy3 += sw + 0.65

# ── Annual Meeting Box ────────────────────────────────────────────────────────
rect(slide, C3X, cy3, C3W, 2.5, LIGHT_GRAY, DARK_GREEN, 0.8)
txt(slide, "2025 Annual Meeting",
    C3X+0.15, cy3+0.12, C3W-0.3, 0.65,
    size=22, bold=True, color=RED, align=PP_ALIGN.CENTER)
txt(slide, "October 7–9, 2025  ·  NC State University, Raleigh, NC",
    C3X+0.15, cy3+0.78, C3W-0.3, 0.38,
    size=10, color=DARK_GRAY, align=PP_ALIGN.CENTER)
txt(slide, "MASBio 5th Annual Meeting:\nBiomass for Value-Added Products — Legacy and Future Pathways",
    C3X+0.15, cy3+1.18, C3W-0.3, 0.65,
    size=9.5, bold=True, color=DARK_GREEN, align=PP_ALIGN.CENTER)

# QR code placeholder
rect(slide, C3X+C3W/2-0.65, cy3+1.9, 1.3, 1.3, WHITE, DARK_GREEN, 1.0)
txt(slide, "QR\nCode", C3X+C3W/2-0.65, cy3+2.25, 1.3, 0.55,
    size=8, bold=True, color=DARK_GREEN, align=PP_ALIGN.CENTER)
cy3 += 2.65

# ── Key Findings ──────────────────────────────────────────────────────────────
section_header(slide, "Key Findings", C3X, cy3, C3W, DARK_GREEN, 12)
cy3 += 0.5
bullet_box(slide, [
    "Automated metadata generation reduces manual effort by ~90%",
    "Hybrid Groq+RAG pipeline runs at zero recurring cost",
    "All MASBIO datasets now findable via Harvard Dataverse search",
    "Ingest-aware publish logic eliminates 'dataset locked' errors",
    "Web UI enables non-technical researchers to publish independently",
    "Agent reasoning handles duplicate detection automatically",
], C3X, cy3, C3W, 2.6, size=9.5, color=DARK_GRAY, bullet_color=BRIGHT_GREEN)
cy3 += 2.7

# ── Acknowledgement ───────────────────────────────────────────────────────────
rect(slide, C3X, cy3, C3W, 1.8, LIGHT_GREEN, MED_GREEN, 0.5)
txt(slide, "Acknowledgements",
    C3X+0.15, cy3+0.1, C3W-0.25, 0.38,
    size=11, bold=True, color=DARK_GREEN)
txt(slide,
    "This work is supported by the Sustainable Agricultural Systems project, "
    "award no. 2020-68012-31881, from the U.S. Department of Agriculture's "
    "National Institute of Food and Agriculture (NIFA). "
    "Data management infrastructure developed as part of MASBIO data stewardship activities.",
    C3X+0.15, cy3+0.5, C3W-0.25, 1.25,
    size=8.5, color=DARK_GRAY)

# ════════════════════════════════════════════════════════════════════════════
# FOOTER
# ════════════════════════════════════════════════════════════════════════════
FY = H_IN - 1.45
rect(slide, 0, FY, W_IN, 1.45, FOOTER_DARK)
rect(slide, 0, FY, W_IN, 0.07, BRIGHT_GREEN)  # green top rule

txt(slide, "Questions?",
    0.3, FY+0.15, 2.5, 0.45,
    size=14, bold=True, color=BRIGHT_GREEN)

txt(slide,
    "Syeda Nyma Ferdous  |  Data Manager\nNorth Carolina State University  |  snferdou@ncsu.edu",
    0.3, FY+0.62, 6.5, 0.75,
    size=9.5, color=WHITE)
txt(slide,
    "Karen Lopez-Olmedo  |  Project Manager\nNC State University  |  kdlopezo@ncsu.edu",
    7.2, FY+0.62, 6.5, 0.75,
    size=9.5, color=WHITE)
txt(slide,
    "Jinxing Wang  |  Project Director\nNC State University  |  jinxing_wang@ncsu.edu",
    14.2, FY+0.62, 6.5, 0.75,
    size=9.5, color=WHITE)

txt(slide, "https://masbio.wvu.edu/home",
    25.5, FY+0.45, 7.2, 0.55,
    size=13, bold=True, color=BRIGHT_GREEN, align=PP_ALIGN.RIGHT)
txt(slide, "dataverse.harvard.edu/dataverse/MASBIO",
    25.5, FY+1.0, 7.2, 0.38,
    size=9, color=rgb("B2DFDB"), align=PP_ALIGN.RIGHT)

txt(slide,
    "This work is supported by the Sustainable Agricultural Systems project, award no. 2020-68012-31881, "
    "from the U.S. Department of Agriculture's National Institute of Food and Agriculture",
    0.3, FY+1.1, 24.0, 0.38,
    size=8.5, color=rgb("AAAAAA"), align=PP_ALIGN.LEFT)

# ── SAVE ─────────────────────────────────────────────────────────────────────
out = Path(__file__).parent / "MASBIO_Poster_2026.pptx"
prs.save(str(out))
print(f"Saved → {out}")

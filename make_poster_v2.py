"""
make_poster_v2.py
-----------------
Reproduces the original MASBio poster visual style with the Data Management
& AI Publishing section added.  Replace coloured placeholder rectangles with
actual photos before printing.

Run:
    pip install python-pptx
    python make_poster_v2.py
Output: MASBIO_Poster_2026_v2.pptx  (33.1 × 27.6 in, ~1.2:1 landscape ratio)
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Slide dimensions ──────────────────────────────────────────────────────────
W, H = 33.1, 27.6   # inches — matches original ~1.2:1 landscape ratio

def rgb(h):
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

# ── Palette (from original poster) ───────────────────────────────────────────
G_DARK   = rgb("1C5C1C")   # MASBio logo / deep headers
G_MED    = rgb("2E7D32")   # section headers
TEAL     = rgb("006666")   # "Our Progress" / timeline
DARK_BG  = rgb("0E1A1A")   # stats panel background
DARK_BG2 = rgb("1A2828")   # annual meeting sub-panel
WHITE    = rgb("FFFFFF")
GRAY_D   = rgb("1A1A1A")
GRAY_M   = rgb("555555")
GRAY_L   = rgb("EEEEEE")
GRAY_LL  = rgb("F7F7F7")
TEAL_ICN = rgb("00827F")   # icon circle fill
BLUE_ICN = rgb("0288D1")   # water drop
GREEN_DM = rgb("76C442")   # data mgmt accent
CYAN_AM  = rgb("00CFFF")   # annual meeting title
FONT     = "Calibri"

# ── Layout constants ──────────────────────────────────────────────────────────
HEAD_H = 4.1       # header height
FOOT_H = 3.1       # footer height
FOOT_Y = H - FOOT_H
CT_Y   = HEAD_H + 0.3          # content area top
CT_BOT = FOOT_Y - 0.25         # content area bottom

# Left research column
LC_X, LC_W = 0.28, 7.9
# Photo / data-management sub-column
PM_X, PM_W = 8.35, 4.0
# Right timeline + stats column
RC_X = 12.6
RC_W = W - RC_X - 0.28

# Dark stats panel starts at this y
DARK_Y = CT_Y + 10.2

# ── Drawing helpers ───────────────────────────────────────────────────────────
def I(v): return Inches(v)
def P(v): return Pt(v)

def box(slide, x, y, w, h, fill, line_clr=None, lw=0.5):
    sh = slide.shapes.add_shape(1, I(x), I(y), I(w), I(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line_clr if line_clr else fill
    if not line_clr:
        sh.line.fill.background()
    else:
        sh.line.width = P(lw)
    return sh

def ov(slide, x, y, d, fill, line_clr=None, lw=1.0):
    sh = slide.shapes.add_shape(9, I(x), I(y), I(d), I(d))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line_clr:
        sh.line.color.rgb = line_clr; sh.line.width = P(lw)
    else:
        sh.line.fill.background()
    return sh

def txt(slide, text, x, y, w, h, sz=10, bold=False, clr=None,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.name = FONT; r.font.size = P(sz)
    r.font.bold = bold; r.font.italic = italic
    if clr: r.font.color.rgb = clr
    return tb

def mtxt(slide, spans, x, y, w, h, base=9, wrap=True):
    """spans = [(text, bold, clr, sz_or_None), ...]"""
    tb = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    for i, (text, bold, clr, sz) in enumerate(spans):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = text
        r.font.name = FONT; r.font.size = P(sz or base)
        r.font.bold = bold
        if clr: r.font.color.rgb = clr
    return tb

def bullets(slide, items, x, y, w, h, sz=9, clr=None, bcl=None):
    clr = clr or GRAY_D
    bcl = bcl or G_MED
    tb = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        rb = p.add_run(); rb.text = "● "
        rb.font.name = FONT; rb.font.size = P(sz)
        rb.font.color.rgb = bcl; rb.font.bold = True
        r = p.add_run(); r.text = item
        r.font.name = FONT; r.font.size = P(sz); r.font.color.rgb = clr

def img_ph(slide, x, y, w, h, label, fill=None, lclr=None, tclr=None):
    fill = fill or rgb("C8E6C9")
    lclr = lclr or G_MED
    tclr = tclr or G_MED
    box(slide, x, y, w, h, fill, lclr, 0.8)
    txt(slide, label, x, y + h/2 - 0.25, w, 0.5,
        sz=9, clr=tclr, align=PP_ALIGN.CENTER, italic=True)

# ═══════════════════════════════════════════════════════════════════════════════
# BUILD SLIDE
# ═══════════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = I(W)
prs.slide_height = I(H)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# White background
box(slide, 0, 0, W, H, WHITE)

# ═══════════════════════════════════════════════════════════════════════════════
# HEADER
# ═══════════════════════════════════════════════════════════════════════════════
# Light green oval behind logo
ov(slide, 0.22, 0.18, 3.0, rgb("D9F0DA"), G_MED, 0.6)
# Leaf shapes (approximated)
ov(slide, 0.50, 0.65, 1.0, G_MED)
ov(slide, 1.15, 0.35, 1.0, G_DARK)
ov(slide, 0.85, 1.05, 0.85, rgb("43A047"))
# Water drop
ov(slide, 1.75, 0.12, 0.65, BLUE_ICN)

# "MASBio" wordmark
txt(slide, "MASBio", 3.4, 0.10, 10.0, 2.6,
    sz=108, bold=True, clr=G_DARK)

# Subtitle line
txt(slide,
    "THE MID-ATLANTIC SUSTAINABLE BIOMASS FOR VALUE-ADDED PRODUCTS CONSORTIUM",
    0.28, 3.35, 16.0, 0.55, sz=11, clr=GRAY_M)

# Thin divider under header
box(slide, 0, HEAD_H - 0.06, W, 0.06, rgb("CCCCCC"))

# "Our Progress" — top right
txt(slide, "Our Progress",
    RC_X + RC_W - 11.5, 0.15, 11.2, 1.5,
    sz=44, bold=True, clr=TEAL, align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════════════════════════════════════════
# LEFT COLUMN — About + Seven Research Areas
# ═══════════════════════════════════════════════════════════════════════════════
cy = CT_Y

# About paragraph
mtxt(slide, [
    ("MASBio", True,  G_DARK, 10.5),
    (" is a consortium of scientists from academia, government agencies and partners "
     "from the wood product industry. We collaborate on integrated and transdisciplinary "
     "research, education and extension project that aims to facilitate and foster development of "
     "the bioeconomy and rural prosperity in the Mid-Atlantic region across seven primary areas:",
     False, GRAY_D, 10),
], LC_X, cy, LC_W, 1.4)
cy += 1.5

# Research areas  (icon_color, header_color, title, description)
AREAS = [
    (TEAL_ICN, G_MED,
     "Feedstock Production",
     "Identify and demonstrate feasible and cost-effective approaches to soil amendment, "
     "vegetation, plant material for bioproduct."),
    (G_MED, G_MED,
     "Harvest and Logistics",
     "Demonstrate efficient and effective harvest and logistics strategies for an "
     "optimized supply chain."),
    (TEAL_ICN, G_MED,
     "Bioproduct Production",
     "Develop and optimize bioproduct conversion processes."),
    (rgb("558B2F"), G_MED,
     "Sustainability Analysis",
     "Evaluate the sustainability and human dimensions of the developed system."),
    (G_DARK, G_MED,
     "System TEA and Optimization",
     "Find the most efficient and economically viable solutions for novel product applications."),
    (TEAL_ICN, G_MED,
     "Education",
     "Engage the next generation of bioproducts leaders through education and internship programs."),
    (G_MED, G_MED,
     "Outreach and Business",
     "Support Outreach and engage with entrepreneurs, stakeholders and business developers "
     "to promote bioeconomic development."),
]

ICON_D = 0.82
ROW_H  = 1.60

for i, (ic, hc, title, desc) in enumerate(AREAS):
    ry = cy + i * ROW_H
    # Circle icon
    ov(slide, LC_X, ry + 0.05, ICON_D, ic)
    # Icon initial letter
    txt(slide, title[0], LC_X, ry + 0.18, ICON_D, ICON_D - 0.35,
        sz=18, bold=True, clr=WHITE, align=PP_ALIGN.CENTER)
    # Header + description
    txt(slide, title,
        LC_X + ICON_D + 0.18, ry, LC_W - ICON_D - 0.22, 0.45,
        sz=10, bold=True, clr=hc)
    txt(slide, desc,
        LC_X + ICON_D + 0.18, ry + 0.46, LC_W - ICON_D - 0.22, 1.1,
        sz=8.5, clr=GRAY_M)

cy += len(AREAS) * ROW_H  # cy now = bottom of research areas

# ═══════════════════════════════════════════════════════════════════════════════
# PHOTO + DATA MANAGEMENT SUB-COLUMN
# ═══════════════════════════════════════════════════════════════════════════════
pm_y = CT_Y

# Photo 1
img_ph(slide, PM_X, pm_y, PM_W, 3.55, "Field Research Photo")
pm_y += 3.65

# Photo 2
img_ph(slide, PM_X, pm_y, PM_W, 3.55, "Team Meeting Photo",
       fill=rgb("DCEEFB"), lclr=BLUE_ICN, tclr=BLUE_ICN)
pm_y += 3.65

# ── DATA MANAGEMENT highlighted box ──────────────────────────────────────────
dm_bot = cy - 0.15            # aligns with bottom of research areas
dm_h   = dm_bot - pm_y
if dm_h < 3.0:                # safety floor
    dm_h = 3.0

DM_BG = rgb("0F2D13")         # very dark green background
box(slide, PM_X, pm_y, PM_W, dm_h, DM_BG)
# Bright left accent stripe
box(slide, PM_X, pm_y, 0.12, dm_h, GREEN_DM)

# Header bar
box(slide, PM_X + 0.12, pm_y, PM_W - 0.12, 0.62, rgb("1B4A20"))
txt(slide, "DATA MANAGEMENT", PM_X + 0.22, pm_y + 0.06, PM_W - 0.26, 0.30,
    sz=9.5, bold=True, clr=WHITE)
txt(slide, "& AI PUBLISHING", PM_X + 0.22, pm_y + 0.35, PM_W - 0.26, 0.25,
    sz=9.5, bold=True, clr=GREEN_DM)

bullets(slide, [
    "Reorganized MASBIO data folder — standardized hierarchy, consistent naming, "
    "consortium-wide access",
    "All datasets published to Harvard Dataverse for permanent open access",
    "Each file assigned a unique, citable DOI — Findable, Accessible, Interoperable, Reusable",
    "AI agent automates full workflow: CSV → metadata → DOI in under 60 seconds",
    "Hybrid Groq LLM + RAG pipeline — zero recurring API cost",
], PM_X + 0.24, pm_y + 0.72, PM_W - 0.32, dm_h - 3.4,
   sz=8.2, clr=WHITE, bcl=GREEN_DM)

# Mini agent flow diagram (compact)
FLOW_Y = pm_y + dm_h - 2.55
flow_items = [
    ("CSV Upload", TEAL_ICN),
    ("RAG + Groq", G_MED),
    ("LangChain Agent", G_DARK),
    ("Harvard Dataverse", rgb("8B0000")),
]
fw = (PM_W - 0.28) / len(flow_items)
for j, (label, fc) in enumerate(flow_items):
    fx = PM_X + 0.22 + j * fw
    box(slide, fx, FLOW_Y, fw - 0.08, 0.52, fc)
    txt(slide, label, fx, FLOW_Y + 0.08, fw - 0.08, 0.36,
        sz=7, bold=True, clr=WHITE, align=PP_ALIGN.CENTER)
    if j < len(flow_items) - 1:
        txt(slide, "→", fx + fw - 0.16, FLOW_Y + 0.1, 0.22, 0.32,
            sz=8, bold=True, clr=GREEN_DM, align=PP_ALIGN.CENTER)

box(slide, PM_X + 0.22, FLOW_Y + 0.58, PM_W - 0.32, 0.06, GREEN_DM)
txt(slide, "✅ DOI returned to researcher in < 60 seconds",
    PM_X + 0.22, FLOW_Y + 0.68, PM_W - 0.32, 0.35,
    sz=7.5, bold=True, clr=GREEN_DM, align=PP_ALIGN.CENTER)

box(slide, PM_X + 0.22, FLOW_Y + 1.08, PM_W - 0.32, 0.40, WHITE,
    rgb("336633"), 0.5)
txt(slide, "🔗  dataverse.harvard.edu  |  DOI: 10.7910/DVN/MASBIO",
    PM_X + 0.22, FLOW_Y + 1.14, PM_W - 0.32, 0.30,
    sz=7.5, bold=True, clr=G_DARK, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# RIGHT COLUMN — Timeline  (above dark panel)
# ═══════════════════════════════════════════════════════════════════════════════
TL_ITEMS = [
    ("2020–2021", "Project started",
     "Collaboration starting with design and layout of field trials, data collection from trials "
     "(soil characteristics and climate and weather data), Biomass/residue samples for analysis."),
    ("2022", "Achievements",
     "Biochar studies, willow/switchgrass trials, and biomass processing exploration. "
     "Soil data, water quality, and emissions analyses. New bioeconomy courses were created. "
     "Biochar demos and business strategies."),
    ("2023", "Expanded outreach",
     "Established new willow trials, optimized bioadhesives and CO2 conversion, "
     "conducted 50+ biomass and soil analyses, economic assessments, and LCA studies. "
     "Expanded outreach with biochar demos and trainings."),
    ("2024", "20+ publications",
     "Advanced biomass and biochar research with 20+ publications, a machine learning model "
     "for soil respiration, and the CarbonCalc biochar credit estimator while engaging 650+ participants. "
     "Developed bioadhesives, SAF, and Char21™ stormwater filtration."),
    ("2025", "5th Annual Meeting",
     "MASBio 5th Annual Meeting: Biomass for Value-Added Products – Legacy and Future Pathways, "
     "October 7-9, 2025, NC State University, Raleigh, NC."),
    ("2026", "Project closure",
     "Data Management & Closure. All datasets archived and published to Harvard Dataverse "
     "with permanent DOIs."),
]

n      = len(TL_ITEMS)
TL_X1  = RC_X + 0.5
TL_X2  = RC_X + RC_W - 0.5
TL_W   = TL_X2 - TL_X1
TL_Y   = CT_Y + 1.4             # horizontal bar y
STEP   = TL_W / (n - 1)
COL_W  = TL_W / n

# Horizontal bar
box(slide, TL_X1, TL_Y + 0.14, TL_W, 0.07, GRAY_M)

for i, (yr, title, desc) in enumerate(TL_ITEMS):
    cx = TL_X1 + i * STEP
    # Year above bar
    txt(slide, yr, cx - 1.0, TL_Y - 0.62, 2.0, 0.42,
        sz=10, bold=True, clr=TEAL, align=PP_ALIGN.CENTER)
    # Circle on bar
    c = slide.shapes.add_shape(9, I(cx - 0.2), I(TL_Y), I(0.40), I(0.40))
    c.fill.solid(); c.fill.fore_color.rgb = TEAL; c.line.fill.background()
    # Description below bar
    mtxt(slide, [
        (title, True,  GRAY_D, 9),
        ("\n" + desc, False, GRAY_M, 8),
    ], cx - COL_W/2 + 0.08, TL_Y + 0.5, COL_W - 0.12, DARK_Y - TL_Y - 0.7)

# ═══════════════════════════════════════════════════════════════════════════════
# DARK STATS PANEL
# ═══════════════════════════════════════════════════════════════════════════════
PANEL_H = FOOT_Y - DARK_Y
box(slide, RC_X, DARK_Y, RC_W, PANEL_H, DARK_BG)

STATS = [
    ("+59 students\ninvolved",         "(Undergraduate, MS, and PhD)"),
    ("+1000 People\nengaged",          ""),
    ("+40 Outreach\nactivities",       ""),
    ("+49 Peer-reviewed\npublications",""),
    ("+200 Academic\nPresentations",   ""),
    ("+15 Programs\ndelivered",        "(College-level courses and teacher training)"),
]

SCOLS = 3
SW    = RC_W / SCOLS
SH    = 3.8
SY    = DARK_Y + 0.45

for i, (num, sub) in enumerate(STATS):
    row = i // SCOLS
    col = i % SCOLS
    sx  = RC_X + col * SW
    sy  = SY + row * SH
    # Icon circle placeholder
    ov(slide, sx + SW/2 - 0.60, sy + 0.1, 1.20, TEAL_ICN)
    txt(slide, "👤" if i in (0,1) else ("📢" if i==2 else ("📄" if i==3 else ("🎓" if i==4 else "🏫"))),
        sx + SW/2 - 0.60, sy + 0.28, 1.20, 0.8,
        sz=16, clr=WHITE, align=PP_ALIGN.CENTER)
    # Stat number
    txt(slide, num, sx + 0.1, sy + 1.45, SW - 0.2, 1.1,
        sz=18, bold=True, clr=WHITE, align=PP_ALIGN.CENTER)
    if sub:
        txt(slide, sub, sx + 0.1, sy + 2.65, SW - 0.2, 0.5,
            sz=7.5, clr=rgb("AAAAAA"), align=PP_ALIGN.CENTER)

# ── Annual Meeting sub-section ─────────────────────────────────────────────────
AM_Y = SY + SH * 2 + 0.5
AM_H = FOOT_Y - AM_Y
box(slide, RC_X, AM_Y, RC_W, AM_H, DARK_BG2)

# Meeting title
txt(slide, "2025 Annual Meeting",
    RC_X + 0.6, AM_Y + 0.25, RC_W - 4.5, 1.6,
    sz=42, bold=True, clr=CYAN_AM, align=PP_ALIGN.CENTER)

txt(slide, "October 7-9, 2025, at NC State University, Raleigh, NC.",
    RC_X + 0.6, AM_Y + 1.95, RC_W - 4.5, 0.55,
    sz=12.5, clr=WHITE, align=PP_ALIGN.CENTER)

txt(slide,
    "Save the date for the MASBio 5th Annual Meeting:\n"
    "Biomass for Value-Added Products – Legacy and Future Pathways",
    RC_X + 0.6, AM_Y + 2.6, RC_W - 4.5, 1.0,
    sz=11, bold=True, clr=WHITE, align=PP_ALIGN.CENTER)

# QR code placeholder (right side)
box(slide, RC_X + RC_W - 3.8, AM_Y + 0.5, 3.2, 3.2, WHITE, DARK_BG, 1.2)
txt(slide, "QR\nCode", RC_X + RC_W - 3.8, AM_Y + 1.6, 3.2, 0.8,
    sz=14, bold=True, clr=DARK_BG, align=PP_ALIGN.CENTER)

# Contact information
txt(slide, "Contact Information",
    RC_X + 0.5, AM_Y + 3.8, RC_W - 4.5, 0.45,
    sz=10.5, bold=True, clr=TEAL, align=PP_ALIGN.CENTER)

CONTACTS = [
    ("Daniel Ciolkosz | Education Leader", "Penn State University", "dec109@psu.edu"),
    ("Karen Lopez-Olmedo | Project Manager", "NC State University", "kdlopezo@ncsu.edu"),
    ("Jinxing Wang | Project Director",  "NC State University", "jingxin_wang@ncsu.edu"),
]
CW = (RC_W - 0.8) / 3
for k, (name, uni, email) in enumerate(CONTACTS):
    cx = RC_X + 0.5 + k * CW
    txt(slide, name,  cx, AM_Y + 4.35, CW - 0.1, 0.35, sz=9.5, bold=True,
        clr=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, uni,   cx, AM_Y + 4.72, CW - 0.1, 0.30, sz=9,   clr=rgb("AAAAAA"),
        align=PP_ALIGN.CENTER)
    txt(slide, email, cx, AM_Y + 5.02, CW - 0.1, 0.30, sz=9,   clr=rgb("00BFFF"),
        align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# FOOTER — Partner Logos
# ═══════════════════════════════════════════════════════════════════════════════
box(slide, 0, FOOT_Y, W, FOOT_H, GRAY_L)
box(slide, 0, FOOT_Y, W, 0.07, rgb("BBBBBB"))  # divider

LOGOS = [
    "West Virginia\nUniversity",  "ESF",  "INL",  "Oak Ridge\nNat. Lab",
    "Penn State",  "CAAFI",  "Univ. of\nGeorgia",  "NC State\nUniversity",
    "USDA",  "Virginia\nTech",  "WV State\nUniversity",  "University\nof Wisconsin",
    "GTI\nEnergy",  "Metzler",  "Dominion\nEnergy",  "ECOSTRAT",
    "AllStar\nEcology",  "TorresSAK",
]
LW = W / len(LOGOS)
for i, logo in enumerate(LOGOS):
    box(slide, i * LW + 0.02, FOOT_Y + 0.12, LW - 0.06, 1.75, WHITE,
        rgb("DDDDDD"), 0.4)
    txt(slide, logo, i * LW + 0.05, FOOT_Y + 0.45, LW - 0.1, 1.1,
        sz=7, clr=GRAY_D, align=PP_ALIGN.CENTER, wrap=True)

# Disclaimer + URL
txt(slide,
    "This work is supported by the Sustainable Agricultural Systems project, "
    "award no. 2020-68012-31881, from the U.S. Department of Agriculture's "
    "National Institute of Food and Agriculture.",
    0.3, FOOT_Y + 2.05, 22.0, 0.75, sz=9, clr=GRAY_M)

txt(slide, "https://masbio.wvu.edu/home",
    23.0, FOOT_Y + 2.1, 9.8, 0.65,
    sz=14, bold=True, clr=G_DARK, align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════════════
out = Path(__file__).parent / "MASBIO_Poster_2026_v2.pptx"
prs.save(str(out))
print(f"Saved  →  {out}")
